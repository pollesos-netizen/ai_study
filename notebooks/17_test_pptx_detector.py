"""
14주차 pptx detector + guide builder 단위 테스트 (TC1~TC21)

13주차 docx detector 테스트 패턴을 따라 작성했습니다.
pptx 특유 케이스(표/그룹/노트/비텍스트 shape)는 TC18~TC21로 추가했습니다.

탐지 함수는 mock으로 주입해 모델 의존성을 끊었습니다.
"""

from __future__ import annotations

import sys
import tempfile
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parent.parent
SRC_DIR = PROJECT_ROOT / "src"
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

from pptx import Presentation
from pptx.util import Inches

from common_apply_result import APPLY_MODE_GUIDE
from deidentify_target_builder import (
    DeidentifyPlan,
    DeidentifyTarget,
    build_deidentify_plan,
)
from pptx_detector import (
    build_guide_for_pptx,
    detect_in_pptx,
    iter_pptx_paragraphs,
    load_pptx,
)


# ── 공통 헬퍼 ──────────────────────────────────────────────────

def _add_textbox_slide(prs, text: str):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    tb.text_frame.text = text
    return slide


def _add_multi_paragraph_textbox(prs, lines: list[str]):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
    tf = tb.text_frame
    if lines:
        tf.text = lines[0]
        for extra in lines[1:]:
            p = tf.add_paragraph()
            p.text = extra
    return slide


def _add_table_slide(prs, cells: list[list[str]]):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    rows = len(cells)
    cols = max(len(r) for r in cells)
    tbl_shape = slide.shapes.add_table(
        rows=rows, cols=cols,
        left=Inches(1), top=Inches(1),
        width=Inches(6), height=Inches(2 * rows),
    )
    tbl = tbl_shape.table
    for r, row_data in enumerate(cells):
        for c, val in enumerate(row_data):
            tbl.cell(r, c).text = val
    return slide


def _make_target(
    *,
    label, matched, start, end, source, action,
    slide_no, section, paragraph_no, context,
    shape_no=None, row_no=None, col_no=None,
    grade="S", order=0,
):
    meta = {
        "fileType": "pptx",
        "slideNo": slide_no,
        "section": section,
        "paragraphNo": paragraph_no,
    }
    if shape_no is not None:
        meta["shapeNo"] = shape_no
    if row_no is not None:
        meta["rowNo"] = row_no
    if col_no is not None:
        meta["colNo"] = col_no

    return DeidentifyTarget(
        label=label, matched=matched, action=action,
        location_label=None,
        location_meta=meta,
        start=start, end=end, source=source,
        reason=f"테스트용 {source} 탐지",
        grade=grade, context=context, order=order,
    )


_runner_path = str(Path(__file__).resolve().parent)
if _runner_path not in sys.path:
    sys.path.insert(0, _runner_path)
from test_helpers import TestRunner, run_test_functions

_runner = TestRunner("14주차 pptx detector 단위 테스트")


def _check(tc_id: str, condition: bool, message: str = "") -> None:
    _runner.check(tc_id, condition, message)


# ── TC1: 이메일이 있는 shape 1개 ──────────────────────────────

def tc1(tmp_dir: Path) -> None:
    print("\nTC1: 이메일이 있는 shape 1개 마스킹")

    prs = Presentation()
    text = "담당자 이메일은 test@example.com입니다."
    _add_textbox_slide(prs, text)
    path = tmp_dir / "tc1.pptx"
    prs.save(str(path))

    matched = "test@example.com"
    start = text.index(matched)
    end = start + len(matched)

    plan = DeidentifyPlan(
        auto_targets=[
            _make_target(
                label="이메일 주소", matched=matched, start=start, end=end,
                source="regex", action="마스킹",
                slide_no=0, section="shape_text", shape_no=0, paragraph_no=0,
                context=text,
            )
        ],
        review_targets=[],
    )

    result = build_guide_for_pptx(str(path), plan)
    _check("TC1.applyMode", result.applyMode == APPLY_MODE_GUIDE)
    _check("TC1.outputFilePath_None", result.outputFilePath is None)
    _check("TC1.autoResults", len(result.autoResults) == 1)
    item = result.autoResults[0]
    _check("TC1.status", item.status == "applied")
    _check("TC1.applied_count", item.appliedTargetCount == 1)
    _check("TC1.preview", "*" * len(matched) in item.appliedText)


# ── TC2: 성명 ──────────────────────────────────────────────────

def tc2(tmp_dir: Path) -> None:
    print("\nTC2: 성명이 있는 shape 1개 마스킹")

    prs = Presentation()
    text = "직원 김도윤의 서류를 검토했습니다."
    _add_textbox_slide(prs, text)
    path = tmp_dir / "tc2.pptx"
    prs.save(str(path))

    s = text.index("김도윤"); e = s + 3
    plan = DeidentifyPlan(
        auto_targets=[
            _make_target(
                label="성명", matched="김도윤", start=s, end=e,
                source="ner", action="마스킹",
                slide_no=0, section="shape_text", shape_no=0, paragraph_no=0,
                context=text,
            )
        ],
        review_targets=[],
    )

    result = build_guide_for_pptx(str(path), plan)
    item = result.autoResults[0]
    _check("TC2.status", item.status == "applied")
    _check("TC2.preview", "***" in item.appliedText)


# ── TC3: 성명 + 이메일 동시 ────────────────────────────────────

def tc3(tmp_dir: Path) -> None:
    print("\nTC3: 한 shape에 성명 + 이메일 동시")

    prs = Presentation()
    text = "담당자 김도윤의 이메일은 test@example.com입니다."
    _add_textbox_slide(prs, text)
    path = tmp_dir / "tc3.pptx"
    prs.save(str(path))

    s1 = text.index("김도윤"); e1 = s1 + 3
    s2 = text.index("test@example.com"); e2 = s2 + len("test@example.com")

    plan = DeidentifyPlan(
        auto_targets=[
            _make_target(label="성명", matched="김도윤", start=s1, end=e1,
                         source="ner", action="마스킹",
                         slide_no=0, section="shape_text", shape_no=0, paragraph_no=0,
                         context=text, order=0),
            _make_target(label="이메일 주소", matched="test@example.com", start=s2, end=e2,
                         source="regex", action="마스킹",
                         slide_no=0, section="shape_text", shape_no=0, paragraph_no=0,
                         context=text, order=1),
        ],
        review_targets=[],
    )

    result = build_guide_for_pptx(str(path), plan)
    _check("TC3.autoResults", len(result.autoResults) == 1)
    item = result.autoResults[0]
    _check("TC3.applied_count", item.appliedTargetCount == 2)
    _check("TC3.label_order", item.label == "성명, 이메일 주소",
           f"label={item.label}")


# ── TC4: 내부 IP 삭제 (delete/mark) ────────────────────────────

def tc4(tmp_dir: Path) -> None:
    print("\nTC4: 내부 IP 삭제 (delete/mark)")

    prs = Presentation()
    text = "서버 IP는 192.168.0.1입니다."
    _add_textbox_slide(prs, text)
    path = tmp_dir / "tc4.pptx"
    prs.save(str(path))

    matched = "192.168.0.1"
    s = text.index(matched); e = s + len(matched)
    plan = DeidentifyPlan(
        auto_targets=[
            _make_target(label="내부 IP 주소", matched=matched, start=s, end=e,
                         source="regex", action="삭제",
                         slide_no=0, section="shape_text", shape_no=0, paragraph_no=0,
                         context=text, grade="C"),
        ],
        review_targets=[],
    )

    result = build_guide_for_pptx(str(path), plan)
    item = result.autoResults[0]
    _check("TC4.delete", matched not in item.appliedText)
    _check("TC4.delete_text", item.appliedText == "서버 IP는 입니다.")

    result_mark = build_guide_for_pptx(str(path), plan, deletion_mode="mark")
    _check("TC4.mark", "(삭제됨)" in result_mark.autoResults[0].appliedText)


# ── TC5: reviewTargets 보존 ────────────────────────────────────

def tc5(tmp_dir: Path) -> None:
    print("\nTC5: reviewTargets 보존")

    prs = Presentation()
    text = "입찰 제안 평가표를 검토했습니다."
    _add_textbox_slide(prs, text)
    path = tmp_dir / "tc5.pptx"
    prs.save(str(path))

    review = _make_target(
        label="민감정보", matched="", start=None, end=None,
        source="ai", action="검토 필요",
        slide_no=0, section="shape_text", shape_no=0, paragraph_no=0,
        context=text, grade="C",
    )

    plan = DeidentifyPlan(auto_targets=[], review_targets=[review])
    result = build_guide_for_pptx(str(path), plan)
    _check("TC5.review_count", len(result.reviewTargets) == 1)
    _check("TC5.auto_empty", len(result.autoResults) == 0)
    rv = result.reviewTargets[0]
    _check("TC5.review_action", rv.action == "검토 필요")


# ── TC6: 필수 필드 누락 (paragraphNo 없음) ──────────────────────

def tc6(tmp_dir: Path) -> None:
    print("\nTC6: paragraphNo 없음 → missing_paragraph_no")

    prs = Presentation()
    _add_textbox_slide(prs, "임의 텍스트.")
    path = tmp_dir / "tc6.pptx"
    prs.save(str(path))

    target = DeidentifyTarget(
        label="이메일", matched="test@example.com", action="마스킹",
        location_label="알 수 없음",
        location_meta={
            "fileType": "pptx",
            "slideNo": 0,
            "section": "shape_text",
            "shapeNo": 0,
            # paragraphNo 누락
        },
        start=0, end=16, source="regex", reason="test",
        grade="S", context="test",
    )
    plan = DeidentifyPlan(auto_targets=[target], review_targets=[])
    result = build_guide_for_pptx(str(path), plan)

    _check("TC6.skipped", result.autoResults[0].status == "skipped")
    _check("TC6.warning",
           any("[missing_paragraph_no]" in w
               for w in result.autoResults[0].warnings))


# ── TC7: 슬라이드 범위 초과 ────────────────────────────────────

def tc7(tmp_dir: Path) -> None:
    print("\nTC7: 슬라이드 범위 초과 → slide_out_of_range")

    prs = Presentation()
    _add_textbox_slide(prs, "단 하나의 슬라이드.")
    path = tmp_dir / "tc7.pptx"
    prs.save(str(path))

    target = _make_target(
        label="성명", matched="홍길동", start=0, end=3,
        source="ner", action="마스킹",
        slide_no=99,  # 범위 초과
        section="shape_text", shape_no=0, paragraph_no=0,
        context="홍길동 무관",
    )
    plan = DeidentifyPlan(auto_targets=[target], review_targets=[])
    result = build_guide_for_pptx(str(path), plan)

    item = result.autoResults[0]
    _check("TC7.skipped", item.status == "skipped")
    _check("TC7.warning",
           any("[slide_out_of_range]" in w for w in item.warnings))


# ── TC8: context 불일치, slice 일치 → 권장 + warning ────────────

def tc8(tmp_dir: Path) -> None:
    print("\nTC8: context 불일치, slice 일치 → 권장 + warning")

    prs = Presentation()
    text = "담당자 이메일은 test@example.com입니다."
    _add_textbox_slide(prs, text)
    path = tmp_dir / "tc8.pptx"
    prs.save(str(path))

    matched = "test@example.com"
    s = text.index(matched); e = s + len(matched)
    target = _make_target(
        label="이메일", matched=matched, start=s, end=e,
        source="regex", action="마스킹",
        slide_no=0, section="shape_text", shape_no=0, paragraph_no=0,
        context="담당자 이메일: test@example.com",  # 실제 paragraph와 다름
    )
    plan = DeidentifyPlan(auto_targets=[target], review_targets=[])
    result = build_guide_for_pptx(str(path), plan)

    item = result.autoResults[0]
    _check("TC8.status_applied", item.status == "applied")
    _check("TC8.applied_count", item.appliedTargetCount == 1)
    _check("TC8.context_mismatch",
           any("[context_mismatch]" in w for w in item.warnings))


# ── TC9: slice 불일치 → skip ───────────────────────────────────

def tc9(tmp_dir: Path) -> None:
    print("\nTC9: slice 불일치 → skip + warning")

    prs = Presentation()
    text = "담당자 이메일은 test@example.com입니다."
    _add_textbox_slide(prs, text)
    path = tmp_dir / "tc9.pptx"
    prs.save(str(path))

    target = _make_target(
        label="이메일", matched="test@example.com",
        start=0, end=5,  # 매치되지 않는 위치
        source="regex", action="마스킹",
        slide_no=0, section="shape_text", shape_no=0, paragraph_no=0,
        context=text,
    )
    plan = DeidentifyPlan(auto_targets=[target], review_targets=[])
    result = build_guide_for_pptx(str(path), plan)

    item = result.autoResults[0]
    _check("TC9.skipped", item.status == "skipped")
    _check("TC9.warning",
           any("[slice_mismatch]" in w for w in item.warnings))


# ── TC10: 여러 슬라이드/shape 분산 ──────────────────────────────

def tc10(tmp_dir: Path) -> None:
    print("\nTC10: 여러 슬라이드/shape에 target 분산")

    prs = Presentation()
    t1 = "이메일: test@example.com"
    t2 = "전화번호: 010-1234-5678"
    _add_textbox_slide(prs, t1)
    _add_textbox_slide(prs, t2)
    path = tmp_dir / "tc10.pptx"
    prs.save(str(path))

    s1 = t1.index("test@example.com"); e1 = s1 + len("test@example.com")
    s2 = t2.index("010-1234-5678"); e2 = s2 + len("010-1234-5678")

    plan = DeidentifyPlan(
        auto_targets=[
            _make_target(label="이메일", matched="test@example.com",
                         start=s1, end=e1, source="regex", action="마스킹",
                         slide_no=0, section="shape_text", shape_no=0, paragraph_no=0,
                         context=t1),
            _make_target(label="전화번호", matched="010-1234-5678",
                         start=s2, end=e2, source="regex", action="마스킹",
                         slide_no=1, section="shape_text", shape_no=0, paragraph_no=0,
                         context=t2),
        ],
        review_targets=[],
    )

    result = build_guide_for_pptx(str(path), plan)
    _check("TC10.count", len(result.autoResults) == 2)
    _check("TC10.all_applied",
           all(it.status == "applied" for it in result.autoResults))


# ── TC11: summary 정합성 ───────────────────────────────────────

def tc11(tmp_dir: Path) -> None:
    print("\nTC11: summary 정합성")

    prs = Presentation()
    t1 = "이메일: test@example.com"
    t2 = "전화번호: 010-1234-5678"
    _add_textbox_slide(prs, t1)
    _add_textbox_slide(prs, t2)
    path = tmp_dir / "tc11.pptx"
    prs.save(str(path))

    s1 = t1.index("test@example.com"); e1 = s1 + len("test@example.com")
    s2 = t2.index("010-1234-5678"); e2 = s2 + len("010-1234-5678")
    plan = DeidentifyPlan(
        auto_targets=[
            _make_target(label="이메일", matched="test@example.com",
                         start=s1, end=e1, source="regex", action="마스킹",
                         slide_no=0, section="shape_text", shape_no=0, paragraph_no=0,
                         context=t1),
            _make_target(label="전화번호", matched="010-1234-5678",
                         start=s2, end=e2, source="regex", action="마스킹",
                         slide_no=1, section="shape_text", shape_no=0, paragraph_no=0,
                         context=t2),
        ],
        review_targets=[],
    )
    result = build_guide_for_pptx(str(path), plan)
    s = result.summary
    _check("TC11.total", s.totalLocations == len(result.autoResults))
    _check("TC11.breakdown",
           s.appliedLocations + s.partialLocations + s.skippedLocations
           == s.totalLocations)
    sum_targets = sum(it.appliedTargetCount + it.skippedTargetCount
                      for it in result.autoResults)
    _check("TC11.autoTargetCount", s.autoTargetCount == sum_targets)


# ── TC12: deletion_mode=mark 표시 ──────────────────────────────

def tc12(tmp_dir: Path) -> None:
    print("\nTC12: deletion_mode=mark")

    prs = Presentation()
    text = "VLAN 100을 사용합니다."
    _add_textbox_slide(prs, text)
    path = tmp_dir / "tc12.pptx"
    prs.save(str(path))

    matched = "VLAN 100"
    s = text.index(matched); e = s + len(matched)
    plan = DeidentifyPlan(
        auto_targets=[
            _make_target(label="VLAN", matched=matched, start=s, end=e,
                         source="regex", action="삭제",
                         slide_no=0, section="shape_text", shape_no=0, paragraph_no=0,
                         context=text, grade="C"),
        ],
        review_targets=[],
    )
    result = build_guide_for_pptx(str(path), plan, deletion_mode="mark")
    _check("TC12.mark", "(삭제됨)" in result.autoResults[0].appliedText)


# ── TC13: 빈 paragraph는 detection 대상 제외 ───────────────────

def tc13(tmp_dir: Path) -> None:
    print("\nTC13: 빈 paragraph는 detection 대상 제외 (paragraphNo 원문 유지)")

    prs = Presentation()
    _add_multi_paragraph_textbox(prs, [
        "",                          # 빈 paragraph (제외, paragraphNo=0)
        "담당자 이메일은 test@example.com입니다.",  # paragraphNo=1
        "",                          # 빈 paragraph (제외, paragraphNo=2)
    ])
    path = tmp_dir / "tc13.pptx"
    prs.save(str(path))

    def mock_regex(text: str):
        if "test@example.com" in text:
            idx = text.index("test@example.com")
            return [{
                "label": "이메일 주소", "value": "test@example.com",
                "start": idx, "end": idx + len("test@example.com"),
                "grade": "S", "action": "마스킹", "desc": "test",
            }]
        return []

    plan = detect_in_pptx(str(path), regex_detect_func=mock_regex)
    _check("TC13.detection_count", len(plan.auto_targets) == 1,
           f"count={len(plan.auto_targets)}")
    _check("TC13.paragraphNo_preserved",
           plan.auto_targets[0].location_meta.get("paragraphNo") == 1,
           f"paragraphNo={plan.auto_targets[0].location_meta.get('paragraphNo')}")


# ── TC14: locationLabel 형식 ───────────────────────────────────

def tc14(tmp_dir: Path) -> None:
    print('\nTC14: locationLabel 형식 검증')

    prs = Presentation()
    long_text = "담당자 김도윤의 이메일은 test@example.com입니다. 길게 이어지는 문장."
    _add_textbox_slide(prs, long_text)
    _add_table_slide(prs, [
        ["이름", "연락처"],
        ["홍길동", "010-1234-5678"],
    ])
    s_notes = prs.slides[0]
    s_notes.notes_slide.notes_text_frame.text = "발표 메모입니다."
    path = tmp_dir / "tc14.pptx"
    prs.save(str(path))

    paragraphs = iter_pptx_paragraphs(load_pptx(str(path)))
    by_section = {p.section: p for p in paragraphs}

    # shape_text: 1번 슬라이드:
    st = by_section.get("shape_text")
    _check("TC14.shape_prefix",
           st.location_label.startswith("1번 슬라이드:"),
           f"label={st.location_label}")
    _check("TC14.shape_truncate", "..." in st.location_label)

    # table_cell: 2번 슬라이드 표 1번 R행 M열:
    tc = by_section.get("table_cell")
    _check("TC14.table_format",
           "2번 슬라이드 표" in tc.location_label
           and "행" in tc.location_label
           and "열" in tc.location_label,
           f"label={tc.location_label}")
    # B안: 항상 "표 N번" 표시 (표가 1개여도 1번이 붙는다)
    _check("TC14.table_index_always_shown",
           "표 1번" in tc.location_label,
           f"label={tc.location_label}")

    # notes: 1번 슬라이드 발표자 노트:
    nt = by_section.get("notes")
    _check("TC14.notes_prefix",
           nt.location_label.startswith("1번 슬라이드 발표자 노트:"),
           f"label={nt.location_label}")


# ── TC15: applyMode + outputFilePath ───────────────────────────

def tc15(tmp_dir: Path) -> None:
    print('\nTC15: applyMode="guide", outputFilePath=None')

    prs = Presentation()
    _add_textbox_slide(prs, "단순 슬라이드입니다.")
    path = tmp_dir / "tc15.pptx"
    prs.save(str(path))

    plan = DeidentifyPlan(auto_targets=[], review_targets=[])
    result = build_guide_for_pptx(str(path), plan)
    _check("TC15.applyMode", result.applyMode == "guide")
    _check("TC15.outputFilePath", result.outputFilePath is None)
    _check("TC15.fileType", result.fileType == "pptx")


# ── TC16: 한 슬라이드의 여러 shape에 target 분산 ───────────────

def tc16(tmp_dir: Path) -> None:
    print("\nTC16: 한 슬라이드의 여러 shape에 target 분산")

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tb1 = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(1))
    tb1.text_frame.text = "이메일: test@example.com"
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(5), Inches(1))
    tb2.text_frame.text = "전화: 010-1234-5678"
    path = tmp_dir / "tc16.pptx"
    prs.save(str(path))

    def mock_regex(text: str):
        results = []
        for needle, label, grade, action in [
            ("test@example.com", "이메일 주소", "S", "마스킹"),
            ("010-1234-5678", "전화번호", "S", "마스킹"),
        ]:
            if needle in text:
                idx = text.index(needle)
                results.append({
                    "label": label, "value": needle,
                    "start": idx, "end": idx + len(needle),
                    "grade": grade, "action": action, "desc": "t",
                })
        return results

    plan = detect_in_pptx(str(path), regex_detect_func=mock_regex)
    _check("TC16.detection_count", len(plan.auto_targets) == 2)

    # 같은 슬라이드의 서로 다른 shape인지 확인
    shape_nos = sorted({t.location_meta.get("shapeNo") for t in plan.auto_targets})
    _check("TC16.different_shapes", len(shape_nos) == 2,
           f"shape_nos={shape_nos}")


# ── TC17: 빈 text_frame은 detection 대상 제외 ──────────────────

def tc17(tmp_dir: Path) -> None:
    print("\nTC17: 빈 text_frame(strip 기준)은 detection 대상 제외")

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tb1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tb1.text_frame.text = "   "  # 공백만
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(5), Inches(1))
    tb2.text_frame.text = "\t\t"  # 탭만
    tb3 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(5), Inches(1))
    tb3.text_frame.text = "실제 텍스트입니다."
    path = tmp_dir / "tc17.pptx"
    prs.save(str(path))

    paragraphs = iter_pptx_paragraphs(load_pptx(str(path)))
    _check("TC17.paragraph_count", len(paragraphs) == 1,
           f"count={len(paragraphs)}")


# ── TC18: 표 shape 내부 셀 탐지 ────────────────────────────────

def tc18(tmp_dir: Path) -> None:
    print("\nTC18: 표 shape 내부 셀 paragraph 탐지")

    prs = Presentation()
    _add_table_slide(prs, [
        ["이름", "이메일"],
        ["홍길동", "hong@example.com"],
    ])
    path = tmp_dir / "tc18.pptx"
    prs.save(str(path))

    def mock_regex(text: str):
        if "@" in text:
            idx = text.find("hong@example.com")
            if idx == -1:
                return []
            return [{
                "label": "이메일 주소", "value": "hong@example.com",
                "start": idx, "end": idx + len("hong@example.com"),
                "grade": "S", "action": "마스킹", "desc": "t",
            }]
        return []

    plan = detect_in_pptx(str(path), regex_detect_func=mock_regex)
    _check("TC18.detected", len(plan.auto_targets) == 1)
    if plan.auto_targets:
        t = plan.auto_targets[0]
        meta = t.location_meta
        _check("TC18.section", meta.get("section") == "table_cell")
        _check("TC18.rowNo", meta.get("rowNo") == 1)
        _check("TC18.colNo", meta.get("colNo") == 1)

    result = build_guide_for_pptx(str(path), plan)
    _check("TC18.guide_applied",
           result.autoResults[0].status == "applied")


# ── TC19: 그룹 shape 내부 텍스트 탐지 ──────────────────────────

def tc19(tmp_dir: Path) -> None:
    print("\nTC19: 그룹 shape 내부 텍스트 탐지 (재귀 순회 검증)")

    # python-pptx는 그룹 생성 API가 직접 제공되지 않으므로
    # XML 조작으로 그룹 shape을 만든다.
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tb1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tb1.text_frame.text = "그룹 안 텍스트: test@example.com"
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(3), Inches(1))
    tb2.text_frame.text = "그룹 안 두 번째 텍스트입니다."

    # 두 텍스트박스를 그룹으로 묶기
    from pptx.oxml.ns import qn
    from copy import deepcopy

    spTree = slide.shapes._spTree
    grpSp = spTree.makeelement(qn("p:grpSp"), {})

    # grpSpPr 추가
    grpSpPr = spTree.makeelement(qn("p:grpSpPr"), {})
    grpSp.append(grpSpPr)

    # nvGrpSpPr 추가 (필수)
    nvGrpSpPr = spTree.makeelement(qn("p:nvGrpSpPr"), {})
    cNvPr = spTree.makeelement(qn("p:cNvPr"),
                               {"id": "100", "name": "Group 1"})
    cNvGrpSpPr = spTree.makeelement(qn("p:cNvGrpSpPr"), {})
    nvSpPr = spTree.makeelement(qn("p:nvPr"), {})
    nvGrpSpPr.append(cNvPr)
    nvGrpSpPr.append(cNvGrpSpPr)
    nvGrpSpPr.append(nvSpPr)
    grpSp.insert(0, nvGrpSpPr)

    # 기존 텍스트박스 두 개를 그룹 안으로 이동
    sp1 = tb1._element
    sp2 = tb2._element
    spTree.remove(sp1)
    spTree.remove(sp2)
    grpSp.append(sp1)
    grpSp.append(sp2)
    spTree.append(grpSp)

    path = tmp_dir / "tc19.pptx"
    prs.save(str(path))

    def mock_regex(text: str):
        if "test@example.com" in text:
            idx = text.index("test@example.com")
            return [{
                "label": "이메일 주소", "value": "test@example.com",
                "start": idx, "end": idx + len("test@example.com"),
                "grade": "S", "action": "마스킹", "desc": "t",
            }]
        return []

    plan = detect_in_pptx(str(path), regex_detect_func=mock_regex)
    # 그룹 안에 있어도 탐지되어야 함
    _check("TC19.detected", len(plan.auto_targets) >= 1,
           f"count={len(plan.auto_targets)}")

    # 전체 paragraph 순회 시에는 그룹 안 두 텍스트박스 모두 잡혀야 함
    paragraphs = iter_pptx_paragraphs(load_pptx(str(path)))
    _check("TC19.group_paragraphs",
           len(paragraphs) == 2,
           f"count={len(paragraphs)}")


# ── TC20: 발표자 노트 탐지 ─────────────────────────────────────

def tc20(tmp_dir: Path) -> None:
    print("\nTC20: 발표자 노트 탐지")

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tb.text_frame.text = "본문 텍스트입니다."
    slide.notes_slide.notes_text_frame.text = "노트: 010-1234-5678"
    path = tmp_dir / "tc20.pptx"
    prs.save(str(path))

    def mock_regex(text: str):
        if "010-1234-5678" in text:
            idx = text.index("010-1234-5678")
            return [{
                "label": "전화번호", "value": "010-1234-5678",
                "start": idx, "end": idx + len("010-1234-5678"),
                "grade": "S", "action": "마스킹", "desc": "t",
            }]
        return []

    plan = detect_in_pptx(str(path), regex_detect_func=mock_regex)
    _check("TC20.detected", len(plan.auto_targets) == 1)
    if plan.auto_targets:
        meta = plan.auto_targets[0].location_meta
        _check("TC20.section_notes", meta.get("section") == "notes")
        _check("TC20.location_label",
               "발표자 노트" in plan.auto_targets[0].location_label,
               f"label={plan.auto_targets[0].location_label}")


# ── TC21: 비텍스트 shape 제외 ──────────────────────────────────

def tc21(tmp_dir: Path) -> None:
    print("\nTC21: 비텍스트 shape(그림 등) 제외")

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tb.text_frame.text = "텍스트 박스"

    # 그림이 없는 환경이라 connector shape으로 대체
    # add_connector로 텍스트가 없는 shape 추가
    from pptx.enum.shapes import MSO_CONNECTOR
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        begin_x=Inches(1), begin_y=Inches(3),
        end_x=Inches(5), end_y=Inches(3),
    )

    path = tmp_dir / "tc21.pptx"
    prs.save(str(path))

    paragraphs = iter_pptx_paragraphs(load_pptx(str(path)))
    # 텍스트 박스 하나만 잡혀야 함 (connector는 has_text_frame=False)
    _check("TC21.text_only",
           len(paragraphs) == 1
           and paragraphs[0].text == "텍스트 박스",
           f"count={len(paragraphs)}")


# ── TC22: 그룹 내부 위치 분리 (key 충돌 방지) ──────────────────

def tc22(tmp_dir: Path) -> None:
    """
    그룹 안의 두 텍스트박스가 같은 paragraphNo=0이어도
    내부적으로 child shape 번호로 구분되어 두 건이 모두 보존되어야 함.
    """
    print("\nTC22: 그룹 내부 위치 분리 (key 충돌 방지)")

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tb1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tb1.text_frame.text = "그룹 안 텍스트 1"
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(3), Inches(1))
    tb2.text_frame.text = "그룹 안 텍스트 2"

    # XML 조작으로 그룹 생성 (TC19와 동일 패턴)
    from pptx.oxml.ns import qn
    spTree = slide.shapes._spTree
    grpSp = spTree.makeelement(qn("p:grpSp"), {})

    nvGrpSpPr = spTree.makeelement(qn("p:nvGrpSpPr"), {})
    nvGrpSpPr.append(spTree.makeelement(qn("p:cNvPr"), {"id": "200", "name": "Group A"}))
    nvGrpSpPr.append(spTree.makeelement(qn("p:cNvGrpSpPr"), {}))
    nvGrpSpPr.append(spTree.makeelement(qn("p:nvPr"), {}))
    grpSp.append(nvGrpSpPr)
    grpSp.append(spTree.makeelement(qn("p:grpSpPr"), {}))

    sp1 = tb1._element
    sp2 = tb2._element
    spTree.remove(sp1)
    spTree.remove(sp2)
    grpSp.append(sp1)
    grpSp.append(sp2)
    spTree.append(grpSp)

    path = tmp_dir / "tc22.pptx"
    prs.save(str(path))

    paragraphs = iter_pptx_paragraphs(load_pptx(str(path)))
    group_paras = [
        p for p in paragraphs
        if p.section == "group_shape_text"
    ]
    _check("TC22.group_paragraphs_count",
           len(group_paras) == 2,
           f"count={len(group_paras)}")

    # 두 paragraph 모두 paragraphNo=0이지만 location key가 달라야 함
    keys = set()
    for p in group_paras:
        # location_key를 직접 만들어보기 (shape_no, group_shape_no, paragraph_no)
        k = (p.section, p.group_shape_no, p.shape_no, p.paragraph_no)
        keys.add(k)
    _check("TC22.keys_unique",
           len(keys) == 2,
           f"keys={keys}")

    # location_label에 "그룹 내부" 표시
    if group_paras:
        _check("TC22.location_label_group_marker",
               "그룹 내부" in group_paras[0].location_label,
               f"label={group_paras[0].location_label}")

    # location_meta에 groupShapeNo 포함
    if group_paras:
        meta = group_paras[0].location_meta
        _check("TC22.location_meta_group_shape_no",
               "groupShapeNo" in meta,
               f"meta={meta}")
        _check("TC22.location_meta_section",
               meta.get("section") == "group_shape_text")

    # build_guide_for_pptx까지 거쳐도 두 항목 모두 보존되는지 확인
    def mock_regex(text: str):
        # 단순히 모든 텍스트에 빈 결과 반환 (탐지 자체는 검증 대상 아님)
        return []

    # 두 paragraph를 가리키는 target을 직접 생성
    targets = []
    for i, p in enumerate(group_paras):
        meta = p.location_meta
        targets.append(DeidentifyTarget(
            label="가짜", matched=p.text[:3] if len(p.text) >= 3 else p.text,
            action="마스킹",
            location_label=p.location_label,
            location_meta=meta,
            start=0, end=min(3, len(p.text)),
            source="regex", reason="test", grade="S",
            context=p.text, order=i,
        ))

    plan = DeidentifyPlan(auto_targets=targets, review_targets=[])
    result = build_guide_for_pptx(str(path), plan)
    # 두 target 모두 자체 항목으로 분리되어야 함 (덮어쓰기 방지)
    _check("TC22.guide_two_items",
           len(result.autoResults) == 2,
           f"count={len(result.autoResults)}")


# ── TC23: 슬라이드 내 여러 표 식별 (tableIndex) ─────────────────

def tc23(tmp_dir: Path) -> None:
    """
    같은 슬라이드에 표가 2개 있을 때 tableIndex로 구분되고
    locationLabel에 "표 N번"이 표시되어야 함.
    """
    print("\nTC23: 같은 슬라이드 내 표 2개 → tableIndex 구분")

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 첫 번째 표
    tbl1_shape = slide.shapes.add_table(
        rows=1, cols=2, left=Inches(1), top=Inches(0.5),
        width=Inches(4), height=Inches(1),
    )
    tbl1 = tbl1_shape.table
    tbl1.cell(0, 0).text = "이름"
    tbl1.cell(0, 1).text = "홍길동"

    # 두 번째 표
    tbl2_shape = slide.shapes.add_table(
        rows=1, cols=2, left=Inches(1), top=Inches(3),
        width=Inches(4), height=Inches(1),
    )
    tbl2 = tbl2_shape.table
    tbl2.cell(0, 0).text = "연락처"
    tbl2.cell(0, 1).text = "010-1234-5678"

    path = tmp_dir / "tc23.pptx"
    prs.save(str(path))

    paragraphs = iter_pptx_paragraphs(load_pptx(str(path)))
    table_paras = [p for p in paragraphs if p.section == "table_cell"]

    # 표 1개당 셀 2개 = 4건
    _check("TC23.table_paragraph_count",
           len(table_paras) == 4,
           f"count={len(table_paras)}")

    # tableIndex가 0, 1로 부여
    table_indices = sorted({p.table_index for p in table_paras})
    _check("TC23.table_indices",
           table_indices == [0, 1],
           f"indices={table_indices}")

    # 첫 번째 표 셀의 locationLabel에 "표 1번"
    first_table_paras = [p for p in table_paras if p.table_index == 0]
    second_table_paras = [p for p in table_paras if p.table_index == 1]

    if first_table_paras:
        _check("TC23.label_table_1",
               "표 1번" in first_table_paras[0].location_label,
               f"label={first_table_paras[0].location_label}")
    if second_table_paras:
        _check("TC23.label_table_2",
               "표 2번" in second_table_paras[0].location_label,
               f"label={second_table_paras[0].location_label}")

    # location_meta에 tableIndex 포함
    if first_table_paras:
        _check("TC23.meta_table_index",
               first_table_paras[0].location_meta.get("tableIndex") == 0)


# ── 실행 ──────────────────────────────────────────────────────

def main() -> None:
    with tempfile.TemporaryDirectory() as tmp_dir_str:
        tmp_dir = Path(tmp_dir_str)
        print("=== 14주차 pptx detector 단위 테스트 ===")
        test_fns = [
            tc1, tc2, tc3, tc4, tc5, tc6, tc7, tc8, tc9,
            tc10, tc11, tc12, tc13, tc14, tc15, tc16, tc17,
            tc18, tc19, tc20, tc21, tc22, tc23,
        ]
        run_test_functions(_runner, test_fns, tmp_dir)

    _runner.report()


if __name__ == "__main__":
    main()
