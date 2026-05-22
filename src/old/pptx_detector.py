"""
pptx 파일 탐지 + 안내(guide) 모드

14주차 1단계: 스켈레톤 + paragraph 순회 함수

목적:
- pptx 파일의 본문 shape, 표 셀, 발표자 노트, 그룹 shape 내부를 paragraph 단위로 순회합니다.
- docx_detector.py와 동일한 패턴을 따릅니다.

처리 범위:
- shape_text: 일반 텍스트 shape (text_frame)
- table_cell: 표 shape의 각 셀
- notes:      발표자 노트
- group:      그룹 shape 한 단계 재귀 분해

범위 외:
- SmartArt, 차트 내부 텍스트
- 그림 캡션, OLE 객체
- 그룹 안의 그룹 (한 단계만 분해)
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable, Iterator

try:
    from src.common_apply_utils import make_location_label_with_context
    from src.deidentify_target_builder import DeidentifyPlan, build_deidentify_plan
except ModuleNotFoundError:
    from common_apply_utils import make_location_label_with_context
    from deidentify_target_builder import DeidentifyPlan, build_deidentify_plan


# ── 데이터 구조 ────────────────────────────────────────────────

@dataclass
class ParsedPptxParagraph:
    """
    pptx의 단일 paragraph를 탐지 단위로 변환한 구조.

    section 값:
    - "shape_text":  shape.text_frame 내부 paragraph
    - "table_cell":  table cell의 text_frame 내부 paragraph
    - "notes":       발표자 노트의 text_frame 내부 paragraph

    paragraphNo 의미 (section별):
    - shape_text:  shape.text_frame.paragraphs 인덱스
    - table_cell:  cell.text_frame.paragraphs 인덱스
    - notes:       notes_text_frame.paragraphs 인덱스

    shapeNo / rowNo / colNo는 section에 따라 사용 여부가 다릅니다.
    """

    slide_no: int  # 슬라이드 인덱스 (0-based)
    section: str  # "shape_text" | "table_cell" | "notes"
    text: str
    paragraph_no: int  # 위 section별 인덱스
    shape_no: int | None = None  # shape_text / table_cell에서 사용
    row_no: int | None = None  # table_cell에서만 사용
    col_no: int | None = None  # table_cell에서만 사용

    @property
    def location_label(self) -> str:
        """
        사용자 표시용 라벨.

        - shape_text: "1번 슬라이드: context..."  (shapeNo는 내부 식별자라 노출 안 함)
        - table_cell: "1번 슬라이드 표 N행 M열: context..."
        - notes:      "1번 슬라이드 발표자 노트: context..."
        """
        slide_idx_1based = self.slide_no + 1

        if self.section == "shape_text":
            base = f"{slide_idx_1based}번 슬라이드"
        elif self.section == "table_cell":
            # row_no/col_no가 None일 가능성은 거의 없지만 방어
            row_disp = (self.row_no + 1) if self.row_no is not None else "?"
            col_disp = (self.col_no + 1) if self.col_no is not None else "?"
            base = f"{slide_idx_1based}번 슬라이드 표 {row_disp}행 {col_disp}열"
        elif self.section == "notes":
            base = f"{slide_idx_1based}번 슬라이드 발표자 노트"
        else:
            base = f"{slide_idx_1based}번 슬라이드 ({self.section})"

        return make_location_label_with_context(base, self.text, max_length=30)

    @property
    def location_meta(self) -> dict[str, Any]:
        meta: dict[str, Any] = {
            "fileType": "pptx",
            "slideNo": self.slide_no,
            "section": self.section,
            "paragraphNo": self.paragraph_no,
        }

        if self.shape_no is not None:
            meta["shapeNo"] = self.shape_no

        if self.row_no is not None:
            meta["rowNo"] = self.row_no

        if self.col_no is not None:
            meta["colNo"] = self.col_no

        return meta


# ── pptx 로드 ──────────────────────────────────────────────────

def load_pptx(input_path: str | Path):
    """
    python-pptx로 pptx 파일을 로드합니다.

    의존성: python-pptx (`pip install python-pptx`)
    """
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError(
            "python-pptx가 설치되어 있지 않습니다. "
            "`pip install python-pptx`를 실행하세요."
        ) from exc

    return Presentation(str(input_path))


# ── shape 순회 보조 함수 ───────────────────────────────────────

def _iter_text_frame_paragraphs(
    text_frame,
) -> Iterator[tuple[int, str]]:
    """
    text_frame의 paragraph를 (paragraph_no, text) 튜플로 yield합니다.

    빈 paragraph(strip 기준)는 건너뜁니다.
    paragraph_no는 빈 paragraph를 포함한 원문 인덱스를 유지합니다.
    """
    for para_index, paragraph in enumerate(text_frame.paragraphs):
        text = paragraph.text

        if not text.strip():
            continue

        yield para_index, text


def _iter_shape_text_paragraphs(
    shape, slide_no: int, shape_no: int,
) -> Iterator[ParsedPptxParagraph]:
    """
    텍스트 shape의 paragraph를 ParsedPptxParagraph로 yield합니다.
    """
    if not shape.has_text_frame:
        return

    for para_index, text in _iter_text_frame_paragraphs(shape.text_frame):
        yield ParsedPptxParagraph(
            slide_no=slide_no,
            section="shape_text",
            text=text,
            paragraph_no=para_index,
            shape_no=shape_no,
        )


def _iter_table_cell_paragraphs(
    shape, slide_no: int, shape_no: int,
) -> Iterator[ParsedPptxParagraph]:
    """
    표 shape의 각 셀 paragraph를 ParsedPptxParagraph로 yield합니다.

    병합 셀 중복 제거는 수행하지 않습니다 (13주차 docx와 동일 정책).
    guide 모드에서는 중복 안내보다 탐지 누락 방지를 우선합니다.
    """
    if not shape.has_table:
        return

    table = shape.table

    for row_index, row in enumerate(table.rows):
        for col_index, cell in enumerate(row.cells):
            # 셀은 항상 text_frame을 가진다고 가정 (python-pptx 보장)
            for para_index, text in _iter_text_frame_paragraphs(cell.text_frame):
                yield ParsedPptxParagraph(
                    slide_no=slide_no,
                    section="table_cell",
                    text=text,
                    paragraph_no=para_index,
                    shape_no=shape_no,
                    row_no=row_index,
                    col_no=col_index,
                )


def _iter_notes_paragraphs(
    slide, slide_no: int,
) -> Iterator[ParsedPptxParagraph]:
    """
    발표자 노트의 paragraph를 ParsedPptxParagraph로 yield합니다.

    notes_slide가 없는 슬라이드는 건너뜁니다.
    """
    if not slide.has_notes_slide:
        return

    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame

    if notes_text_frame is None:
        return

    for para_index, text in _iter_text_frame_paragraphs(notes_text_frame):
        yield ParsedPptxParagraph(
            slide_no=slide_no,
            section="notes",
            text=text,
            paragraph_no=para_index,
        )


def _iter_shape_recursive(
    shape, slide_no: int, shape_no: int,
) -> Iterator[ParsedPptxParagraph]:
    """
    그룹 shape을 한 단계 재귀 분해하여 내부의 텍스트/표 shape을 순회합니다.

    그룹 안의 또 다른 그룹은 14주차 PoC 범위 외이므로 무시합니다.
    """
    # python-pptx의 GROUP shape_type은 MSO_SHAPE_TYPE.GROUP (6)
    # has_table / has_text_frame은 그룹 자체에는 False이므로
    # shape.shapes로 내부 shape에 접근합니다.

    if not hasattr(shape, "shapes"):
        return

    for child_index, child in enumerate(shape.shapes):
        # 그룹 안의 그룹은 14주차 범위 외
        if hasattr(child, "shapes"):
            # 한 단계만 들어가고, 그 안의 그룹은 무시
            continue

        if child.has_table:
            yield from _iter_table_cell_paragraphs(child, slide_no, shape_no)
        elif child.has_text_frame:
            yield from _iter_shape_text_paragraphs(child, slide_no, shape_no)
        # 그림/차트 등은 무시


def _is_group_shape(shape) -> bool:
    """
    그룹 shape 판단.

    python-pptx에서 그룹 shape은 shape.shapes 속성을 가집니다.
    한편 has_table=False, has_text_frame=False입니다.
    """
    # MSO_SHAPE_TYPE.GROUP == 6
    try:
        return shape.shape_type == 6
    except (AttributeError, ValueError):
        return hasattr(shape, "shapes") and not shape.has_text_frame and not shape.has_table


# ── 전체 순회 ──────────────────────────────────────────────────

def iter_pptx_paragraphs(prs) -> list[ParsedPptxParagraph]:
    """
    pptx 파일의 모든 paragraph를 탐지 단위로 반환합니다.

    순회 대상:
    - 본문 shape (text_frame)
    - 표 shape (table cell)
    - 발표자 노트
    - 그룹 shape 한 단계 재귀

    순회 제외:
    - 빈 paragraph (strip 기준)
    - SmartArt, 차트, 그림, OLE 객체
    - 그룹 안의 그룹
    """
    results: list[ParsedPptxParagraph] = []

    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if _is_group_shape(shape):
                results.extend(
                    _iter_shape_recursive(shape, slide_index, shape_index)
                )
                continue

            if shape.has_table:
                results.extend(
                    _iter_table_cell_paragraphs(shape, slide_index, shape_index)
                )
                continue

            if shape.has_text_frame:
                results.extend(
                    _iter_shape_text_paragraphs(shape, slide_index, shape_index)
                )
                continue

            # 그 외 (그림/차트/SmartArt/OLE 등)은 14주차 범위 외

        # 발표자 노트는 슬라이드 단위
        results.extend(_iter_notes_paragraphs(slide, slide_index))

    return results


# ── Detection 어댑터 ──────────────────────────────────────────

def _make_target_dict_from_regex(
    raw: Any,
    paragraph: ParsedPptxParagraph,
    order: int,
) -> dict[str, Any] | None:
    """
    regex_detector의 결과(DetectionResult 또는 유사 dict)를 Detection dict로 변환합니다.

    docx_detector._make_target_dict_from_regex와 동일한 패턴.
    """
    def _get(obj, *names, default=None):
        for name in names:
            if isinstance(obj, dict) and name in obj:
                return obj.get(name)
            if hasattr(obj, name):
                return getattr(obj, name)
        return default

    label = _get(raw, "label")
    value = _get(raw, "value", "matched", "match")
    start = _get(raw, "start")
    end = _get(raw, "end")
    grade = _get(raw, "grade", default="S")
    action = _get(raw, "action", default="마스킹")
    desc = _get(raw, "desc", "reason", default=None)

    if label is None or value is None or start is None or end is None:
        return None

    return {
        "label": str(label),
        "matched": str(value),
        "grade": str(grade),
        "action": str(action),
        "source": "regex",
        "context": paragraph.text,
        "locationLabel": paragraph.location_label,
        "locationMeta": paragraph.location_meta,
        "start": int(start),
        "end": int(end),
        "sensitiveType": _get(raw, "sensitive_type", "sensitiveType", default=None),
        "sensitiveCategory": _get(raw, "sensitive_category", "sensitiveCategory", default=label),
        "reason": str(desc) if desc else f"정규식 탐지: {label}",
        "_order": order,
    }


def _make_target_dict_from_ner(
    raw: dict[str, Any],
    paragraph: ParsedPptxParagraph,
    order: int,
    *,
    threshold: float,
) -> dict[str, Any] | None:
    """
    Hugging Face NER 출력(aggregation_strategy="simple" 기준)을 Detection dict로 변환합니다.

    PERSON 계열 라벨만 변환합니다.
    confidence < threshold이면 None을 반환합니다.
    """
    entity_label = (raw.get("entity_group") or raw.get("entity") or "").upper()
    entity_label = entity_label.replace("B-", "").replace("I-", "")

    if entity_label not in {"PERSON", "PER", "PS", "인명"}:
        return None

    score = float(raw.get("score") or 0.0)
    if score < threshold:
        return None

    start = raw.get("start")
    end = raw.get("end")

    if start is None or end is None:
        return None

    start = int(start)
    end = int(end)

    matched = paragraph.text[start:end] or raw.get("word") or ""

    return {
        "label": "성명",
        "matched": str(matched),
        "grade": "S",
        "action": "마스킹",
        "source": "ner",
        "context": paragraph.text,
        "locationLabel": paragraph.location_label,
        "locationMeta": paragraph.location_meta,
        "start": start,
        "end": end,
        "sensitiveType": "개인정보",
        "sensitiveCategory": "성명",
        "reason": (
            f"NER 모델 PERSON 탐지 / original_label={raw.get('entity_group') or raw.get('entity')}"
            f" / confidence={score:.4f} / threshold={threshold:.2f}"
        ),
        "_order": order,
    }


def _make_target_dict_from_ai(
    grade: str,
    confidence: float,
    paragraph: ParsedPptxParagraph,
    order: int,
    *,
    threshold: float,
    prob_map: dict[str, float] | None = None,
) -> dict[str, Any] | None:
    """
    AI 문장분류 결과를 review target dict로 변환합니다.

    AI Detection은 start/end가 없고 matched가 빈 문자열입니다.
    grade=='O'이거나 confidence < threshold이면 review target을 만들지 않습니다.
    """
    if grade == "O" or confidence < threshold:
        return None

    prob_text = ""
    if prob_map:
        prob_text = " / probs=(" + ", ".join(
            f"{label}={prob:.3f}" for label, prob in prob_map.items()
        ) + ")"

    return {
        "label": "민감정보",
        "matched": "",
        "grade": grade,
        "action": "검토 필요",
        "source": "ai",
        "context": paragraph.text,
        "locationLabel": paragraph.location_label,
        "locationMeta": paragraph.location_meta,
        "start": None,
        "end": None,
        "sensitiveType": "문맥 기반 민감정보",
        "sensitiveCategory": f"AI_{grade}",
        "reason": (
            f"AI 문장분류 grade={grade} / confidence={confidence:.4f}"
            f" / threshold={threshold:.2f}{prob_text}"
        ),
        "_order": order,
    }


# ── 탐지 파이프라인 ────────────────────────────────────────────

def detect_in_pptx(
    input_path: str,
    *,
    regex_detect_func: Callable[[str], list[Any]] | None = None,
    ner_detect_func: Callable[[str], list[dict[str, Any]]] | None = None,
    ai_predict_func: Callable[[str], tuple[str, float, dict[str, float]]] | None = None,
    ner_threshold: float = 0.8,
    ai_threshold: float = 0.6,
) -> DeidentifyPlan:
    """
    pptx 파일의 모든 paragraph(본문 shape / 표 셀 / 발표자 노트 / 그룹 내부)를
    순회하며 탐지를 수행하고 DeidentifyPlan을 생성합니다.

    Args:
        input_path: pptx 파일 경로
        regex_detect_func: text -> regex detection 목록 (생략 시 regex_detector.detect_patterns 사용)
        ner_detect_func: text -> HF NER pipeline 출력 목록 (생략 시 NER skip)
        ai_predict_func: text -> (grade, confidence, prob_map) (생략 시 AI skip)
        ner_threshold: NER confidence 임계값
        ai_threshold: AI confidence 임계값

    Returns:
        DeidentifyPlan (auto_targets + review_targets)

    탐지 함수를 주입형으로 받는 이유:
    - 단위 테스트에서 모델 의존성을 끊을 수 있도록 합니다.
    - 14주차 초반에는 regex만 연결해서 guide 구조부터 검증할 수 있습니다.
    """
    # regex 탐지 함수 기본값
    if regex_detect_func is None:
        try:
            from src.regex_detector import detect_patterns as _detect_patterns
        except ModuleNotFoundError:
            from regex_detector import detect_patterns as _detect_patterns
        regex_detect_func = _detect_patterns

    prs = load_pptx(input_path)
    paragraphs = iter_pptx_paragraphs(prs)

    detections: list[dict[str, Any]] = []
    order = 0

    for paragraph in paragraphs:
        # regex
        raw_regex = regex_detect_func(paragraph.text) or []
        for raw in raw_regex:
            detection = _make_target_dict_from_regex(raw, paragraph, order)
            if detection is not None:
                detections.append(detection)
                order += 1

        # NER
        if ner_detect_func is not None:
            try:
                raw_ner = ner_detect_func(paragraph.text) or []
            except Exception as exc:
                print(f"[NER] {paragraph.location_label} 탐지 실패: {exc}")
                raw_ner = []

            for raw in raw_ner:
                detection = _make_target_dict_from_ner(
                    raw, paragraph, order, threshold=ner_threshold,
                )
                if detection is not None:
                    detections.append(detection)
                    order += 1

        # AI
        if ai_predict_func is not None:
            try:
                grade, confidence, prob_map = ai_predict_func(paragraph.text)
            except Exception as exc:
                print(f"[AI] {paragraph.location_label} 예측 실패: {exc}")
                grade, confidence, prob_map = "O", 0.0, {}

            if grade is not None and confidence is not None:
                detection = _make_target_dict_from_ai(
                    grade, confidence, paragraph, order,
                    threshold=ai_threshold, prob_map=prob_map,
                )
                if detection is not None:
                    detections.append(detection)
                    order += 1

    return build_deidentify_plan(detections)


# ── guide 생성 ─────────────────────────────────────────────────

# pptx 전용 warning type은 common_apply_utils.py로 이전했습니다.
try:
    from src.common_apply_utils import (
        WARNING_MISSING_SLIDE_NO,
        WARNING_SLIDE_OUT_OF_RANGE,
        WARNING_SHAPE_NOT_FOUND,
        WARNING_UNKNOWN_SECTION,
    )
except ModuleNotFoundError:
    from common_apply_utils import (
        WARNING_MISSING_SLIDE_NO,
        WARNING_SLIDE_OUT_OF_RANGE,
        WARNING_SHAPE_NOT_FOUND,
        WARNING_UNKNOWN_SECTION,
    )


def _make_pptx_location_key(meta: dict[str, Any]) -> tuple | None:
    """
    pptx target의 위치 그룹화 키를 생성합니다.

    section 값에 따라 키 구성이 다릅니다.
    필수 필드가 누락되면 None을 반환합니다.

    - shape_text:  (slideNo, "shape_text", shapeNo, paragraphNo)
    - table_cell:  (slideNo, "table_cell", shapeNo, rowNo, colNo, paragraphNo)
    - notes:       (slideNo, "notes", paragraphNo)
    """
    slide_no = meta.get("slideNo")
    section = str(meta.get("section") or "")
    paragraph_no = meta.get("paragraphNo")

    if slide_no is None or paragraph_no is None:
        return None

    if section == "shape_text":
        shape_no = meta.get("shapeNo")
        if shape_no is None:
            return None
        return (int(slide_no), "shape_text", int(shape_no), int(paragraph_no))

    if section == "table_cell":
        shape_no = meta.get("shapeNo")
        row_no = meta.get("rowNo")
        col_no = meta.get("colNo")
        if shape_no is None or row_no is None or col_no is None:
            return None
        return (
            int(slide_no), "table_cell", int(shape_no),
            int(row_no), int(col_no), int(paragraph_no),
        )

    if section == "notes":
        return (int(slide_no), "notes", int(paragraph_no))

    # 알 수 없는 section
    return None


def _format_location_meta_for_label(meta: dict[str, Any]) -> str:
    """
    경고 메시지/locationLabel 보조용 간단 표현.
    """
    slide_disp = (
        f"{int(meta['slideNo']) + 1}번 슬라이드"
        if meta.get("slideNo") is not None else "슬라이드?"
    )
    section = meta.get("section")

    if section == "shape_text":
        return slide_disp
    if section == "table_cell":
        row = meta.get("rowNo")
        col = meta.get("colNo")
        row_disp = (row + 1) if row is not None else "?"
        col_disp = (col + 1) if col is not None else "?"
        return f"{slide_disp} 표 {row_disp}행 {col_disp}열"
    if section == "notes":
        return f"{slide_disp} 발표자 노트"
    return f"{slide_disp} ({section})"


def _index_pptx_paragraphs(
    prs,
) -> dict[tuple, str]:
    """
    pptx에서 paragraph 위치 키 → text 매핑을 미리 생성합니다.

    빈 paragraph(strip 기준)는 포함하지 않으므로 build_guide에서
    그런 key는 자동으로 "범위 외" 처리됩니다.
    """
    index: dict[tuple, str] = {}

    for paragraph in iter_pptx_paragraphs(prs):
        meta = paragraph.location_meta
        key = _make_pptx_location_key(meta)
        if key is None:
            continue
        index[key] = paragraph.text

    return index


def _group_targets_by_location(
    targets,
):
    """
    auto target을 pptx location key 기준으로 묶습니다.

    fileType이 pptx가 아닌 target은 제외합니다.
    필수 필드가 누락되거나 알 수 없는 section이면 skipped item으로 처리합니다.
    """
    try:
        from src.common_apply_result import CommonApplyItem
        from src.common_apply_utils import (
            WARNING_MISSING_PARAGRAPH_NO,
            format_warning,
        )
    except ModuleNotFoundError:
        from common_apply_result import CommonApplyItem
        from common_apply_utils import (
            WARNING_MISSING_PARAGRAPH_NO,
            format_warning,
        )

    grouped: dict[tuple, list] = {}
    skipped_items: list = []
    warnings: list[str] = []

    for target in targets:
        meta = target.location_meta or {}

        if str(meta.get("fileType") or "").lower() != "pptx":
            continue

        section = str(meta.get("section") or "")
        key = _make_pptx_location_key(meta)

        if key is None:
            # 어떤 필드가 빠졌는지 진단
            label = target.location_label or _format_location_meta_for_label(meta)

            if section not in {"shape_text", "table_cell", "notes"}:
                warning = format_warning(
                    WARNING_UNKNOWN_SECTION,
                    f"{label}: 알 수 없는 section={section!r}이므로 안내를 생성하지 않습니다.",
                )
            elif meta.get("slideNo") is None:
                warning = format_warning(
                    WARNING_MISSING_SLIDE_NO,
                    f"{label}: slideNo가 없어 안내를 생성하지 못했습니다.",
                )
            elif meta.get("paragraphNo") is None:
                warning = format_warning(
                    WARNING_MISSING_PARAGRAPH_NO,
                    f"{label}: paragraphNo가 없어 안내를 생성하지 못했습니다.",
                )
            else:
                # shape_text/table_cell의 보조 필드 누락
                warning = format_warning(
                    WARNING_MISSING_PARAGRAPH_NO,
                    f"{label}: 위치 필드(shapeNo/rowNo/colNo)가 누락되어 안내를 생성하지 못했습니다.",
                )

            warnings.append(warning)
            skipped_items.append(
                CommonApplyItem(
                    locationLabel=target.location_label or label,
                    locationMeta=meta,
                    label=target.label or "",
                    action=target.action,
                    originalText=target.context or "",
                    appliedText=target.context or "",
                    status="skipped",
                    appliedTargetCount=0,
                    skippedTargetCount=1,
                    warnings=[warning],
                )
            )
            continue

        grouped.setdefault(key, []).append(target)

    return grouped, skipped_items, warnings


def _build_guide_item_for_location(
    key: tuple,
    paragraph_text: str | None,  # None이면 범위 초과 / 빈 paragraph
    targets,
    *,
    deletion_mode: str,
):
    """
    한 location에 속한 target 목록에 대해 guide 모드 CommonApplyItem을 생성합니다.
    """
    try:
        from src.common_apply_result import CommonApplyItem
        from src.common_apply_utils import (
            WARNING_CONTEXT_MISMATCH,
            WARNING_EMPTY_PARAGRAPH_TARGET,
            actions_for_targets,
            format_warning,
            labels_for_targets,
            make_status,
            normalize_nfc,
            validate_slice_against_text,
        )
        from src.deidentify_apply import apply_targets_to_text
    except ModuleNotFoundError:
        from common_apply_result import CommonApplyItem
        from common_apply_utils import (
            WARNING_CONTEXT_MISMATCH,
            WARNING_EMPTY_PARAGRAPH_TARGET,
            actions_for_targets,
            format_warning,
            labels_for_targets,
            make_status,
            normalize_nfc,
            validate_slice_against_text,
        )
        from deidentify_apply import apply_targets_to_text

    representative = targets[0]
    location_label = representative.location_label or _format_location_meta_for_label(
        representative.location_meta or {}
    )
    location_meta = representative.location_meta or {}

    warnings: list[str] = []

    # paragraph_text가 None이면 범위 초과 (탐지 시점 이후 슬라이드/shape이 사라졌거나 잘못된 위치)
    if paragraph_text is None:
        warning = format_warning(
            WARNING_SLIDE_OUT_OF_RANGE,
            f"{location_label}: 위치가 현재 pptx에서 발견되지 않습니다 (key={key}).",
        )
        warnings.append(warning)
        return CommonApplyItem(
            locationLabel=location_label,
            locationMeta=location_meta,
            label=labels_for_targets(targets),
            action=actions_for_targets(targets),
            originalText="",
            appliedText="",
            status="skipped",
            appliedTargetCount=0,
            skippedTargetCount=len(targets),
            warnings=warnings,
        )

    # 빈 paragraph 방어
    if not paragraph_text.strip():
        warning = format_warning(
            WARNING_EMPTY_PARAGRAPH_TARGET,
            f"{location_label}: 빈 paragraph를 가리키는 target은 안내를 생성하지 않습니다.",
        )
        warnings.append(warning)
        return CommonApplyItem(
            locationLabel=location_label,
            locationMeta=location_meta,
            label=labels_for_targets(targets),
            action=actions_for_targets(targets),
            originalText="",
            appliedText="",
            status="skipped",
            appliedTargetCount=0,
            skippedTargetCount=len(targets),
            warnings=warnings,
        )

    # context 불일치 (적용은 진행)
    if any(
        target.context is not None
        and normalize_nfc(target.context) != normalize_nfc(paragraph_text)
        for target in targets
    ):
        warnings.append(
            format_warning(
                WARNING_CONTEXT_MISMATCH,
                f"{location_label}: target.context와 실제 paragraph 텍스트가 다릅니다. "
                "paragraph 텍스트 기준으로 slice 검증 후 권장 여부를 판단합니다.",
            )
        )

    # slice 검증
    valid_targets = []
    skipped_count = 0

    for target in targets:
        warning_type, slice_error = validate_slice_against_text(paragraph_text, target)
        if slice_error is not None:
            warnings.append(
                format_warning(warning_type, f"{location_label}: {slice_error}")
            )
            skipped_count += 1
            continue
        valid_targets.append(target)

    # guide 모드 preview 생성 (메모리에서만)
    if valid_targets:
        apply_result = apply_targets_to_text(
            paragraph_text,
            valid_targets,
            deletion_mode=deletion_mode,
        )
        applied_text = apply_result.applied_text
        warnings.extend(apply_result.warnings)
        applied_count = len(apply_result.applied_targets)
        skipped_count += len(apply_result.skipped_targets)
    else:
        applied_text = paragraph_text
        applied_count = 0

    return CommonApplyItem(
        locationLabel=location_label,
        locationMeta=location_meta,
        label=labels_for_targets(targets),
        action=actions_for_targets(targets),
        originalText=paragraph_text,
        appliedText=applied_text,
        status=make_status(applied_count, skipped_count),
        appliedTargetCount=applied_count,
        skippedTargetCount=skipped_count,
        warnings=warnings,
    )


def build_guide_for_pptx(
    input_path: str,
    plan: DeidentifyPlan,
    *,
    deletion_mode: str = "delete",
):
    """
    DeidentifyPlan을 받아 guide 모드 CommonApplyResult를 생성합니다.

    실제 파일을 수정하지 않으며, outputFilePath는 None입니다.
    """
    try:
        from src.common_apply_result import (
            APPLY_MODE_GUIDE,
            CommonApplyResult,
            build_summary,
            make_review_items,
        )
    except ModuleNotFoundError:
        from common_apply_result import (
            APPLY_MODE_GUIDE,
            CommonApplyResult,
            build_summary,
            make_review_items,
        )

    prs = load_pptx(input_path)
    paragraph_index = _index_pptx_paragraphs(prs)

    grouped, skipped_items, global_warnings = _group_targets_by_location(plan.auto_targets)

    auto_results = list(skipped_items)

    for key, targets in grouped.items():
        paragraph_text = paragraph_index.get(key)  # 없으면 None → 범위 외 처리
        item = _build_guide_item_for_location(
            key,
            paragraph_text,
            targets,
            deletion_mode=deletion_mode,
        )
        auto_results.append(item)

    review_items = make_review_items(plan.review_targets)
    summary = build_summary(auto_results, review_items, global_warnings)

    return CommonApplyResult(
        fileType="pptx",
        applyMode=APPLY_MODE_GUIDE,
        inputFilePath=str(input_path),
        outputFilePath=None,
        autoResults=auto_results,
        reviewTargets=review_items,
        warnings=global_warnings,
        summary=summary,
    )


def detect_and_build_guide_for_pptx(
    input_path: str,
    *,
    regex_detect_func: Callable[[str], list[Any]] | None = None,
    ner_detect_func: Callable[[str], list[dict[str, Any]]] | None = None,
    ai_predict_func: Callable[[str], tuple[str, float, dict[str, float]]] | None = None,
    ner_threshold: float = 0.8,
    ai_threshold: float = 0.6,
    deletion_mode: str = "delete",
):
    """
    detect_in_pptx + build_guide_for_pptx 편의 wrapper.
    """
    plan = detect_in_pptx(
        input_path,
        regex_detect_func=regex_detect_func,
        ner_detect_func=ner_detect_func,
        ai_predict_func=ai_predict_func,
        ner_threshold=ner_threshold,
        ai_threshold=ai_threshold,
    )

    return build_guide_for_pptx(
        input_path,
        plan,
        deletion_mode=deletion_mode,
    )
